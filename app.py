from flask import Flask, render_template, request, send_file
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import os
import cv2
import numpy as np
from PIL import Image
from io import BytesIO

app = Flask(__name__)

RAPIDAPI_URL = "https://google-search-master-mega.p.rapidapi.com/shopping"
HEADERS = {
    "x-rapidapi-host": "google-search-master-mega.p.rapidapi.com",
    "x-rapidapi-key": "3ff8754e1bmsh6199b185d80c9c9p14c709jsn22be61e4d3f3"  # Replace with your key
}

# Ensure static folder exists
STATIC_FOLDER = os.path.join(os.path.dirname(__file__), 'static')
if not os.path.exists(STATIC_FOLDER):
    os.makedirs(STATIC_FOLDER)

# Function to download image in the highest quality
def download_high_quality_image(image_url):
    try:
        # Get image as bytes
        response = requests.get(image_url, stream=True, timeout=10)
        if response.status_code == 200:
            image_data = response.content
            return image_data
        else:
            raise Exception(f"Error fetching image: {response.status_code}")
    except Exception as e:
        print(f"Error downloading image: {e}")
        return None

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        query = request.form["product"]
        params = {
            "q": query,
            "gl": "us",
            "hl": "en",
            "autocorrect": "true",
            "num": "10",
            "page": "1"
        }

        response = requests.get(RAPIDAPI_URL, headers=HEADERS, params=params)
        if response.status_code == 200:
            data = response.json()
            if "shopping" in data:
                products = [{
                    "name": item.get("title", "No Title"),
                    "image": item.get("imageUrl", ""),
                    "url": item.get("link", "")
                } for item in data["shopping"]]
                return render_template("results.html", products=products, query=query)
            else:
                return f"No shopping results found in response: {data}"
        else:
            return f"API Error: {response.status_code} - {response.text}"
    return render_template("index.html")


@app.route("/generate_ppt", methods=["POST"])
def generate_ppt():
    selected = request.form.getlist("selected_products")
    names = request.form.getlist("name")
    images = request.form.getlist("image")
    urls = request.form.getlist("url")

    prs = Presentation()
    for idx in selected:
        idx = int(idx)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = names[idx]

        # Download and add high-quality image
        try:
            image_data = download_high_quality_image(images[idx])
            if image_data:
                # Save image as PNG or JPG in highest resolution
                image = Image.open(BytesIO(image_data))
                image_path = os.path.join(STATIC_FOLDER, f"product_image_{idx}.png")
                image.save(image_path)

                # Add image to PowerPoint
                slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), height=Inches(3.5))

                # Cleanup the image after adding
                os.remove(image_path)
            else:
                print(f"Failed to download image for product {names[idx]}")

        except Exception as e:
            print(f"Error processing image for {names[idx]}: {e}")

        # Add product URL directly (no extra text)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = urls[idx]  # Directly save the product URL
        run.font.size = Pt(12)
        run.hyperlink.address = urls[idx]

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)

    return send_file(ppt_stream, as_attachment=True, download_name="Selected_Products.pptx")


if __name__ == "__main__":
    app.run(debug=True)
